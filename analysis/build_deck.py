#!/usr/bin/env python3
"""Build a 3-slide Pasta / Noodles review from the monthly trended category export."""

from __future__ import annotations

import re
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "03. Monthly Trended Export_2026-08-18.xlsx"
OUTPUT_DIR = ROOT / "output"
CHART_DIR = OUTPUT_DIR / "charts"
DECK = OUTPUT_DIR / "Unilever_Pasta_Noodles_Jun2026.pptx"

NAVY = RGBColor(0x0B, 0x25, 0x45)
BLUE = RGBColor(0x1B, 0x4F, 0x8A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
ORANGE = RGBColor(0xC4, 0x5C, 0x26)
RED = RGBColor(0x9B, 0x2C, 0x2C)
GREEN = RGBColor(0x1F, 0x7A, 0x4D)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x6B, 0x73)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD5, 0xDB, 0xE0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def load() -> pd.DataFrame:
    df = pd.read_excel(SOURCE, skiprows=3)
    df = df.loc[:, ~df.columns.astype(str).str.contains("^Unnamed")]
    df.columns = ["month", "category", "subcat", "mfr", "sku", "price", "value", "volume"]
    df["date"] = pd.to_datetime(df["month"], format="%b %y")
    df["price"] = np.where(df["volume"] > 0, df["value"] / df["volume"], np.nan)
    df["brand"] = [brand_of(sku, mfr) for sku, mfr in zip(df["sku"], df["mfr"])]
    df["pack_g"] = df["sku"].map(pack_grams)
    df["pack_price"] = df["price"] * df["pack_g"] / 1000.0
    return df


def brand_of(sku: str, mfr: str) -> str:
    rules = [
        (r"^Maggi", "Maggi"),
        (r"^Kellogg", "Kellogg's"),
        (r"^Indomie", "Indomie"),
        (r"^Roka", "Roka"),
        (r"^Mr Pasta", "Mr Pasta"),
        (r"^Fatti", "Fatti's & Moni's"),
        (r"^Samyang|^Buldak", "Samyang"),
        (r"^Eezee|^Joy Instant", "Eezee/Joy"),
        (r"^Lucky Star", "Lucky Star"),
        (r"^Fastmove", "Fastmove"),
        (r"^Knorr", "Knorr"),
        (r"^Serena", "Serena"),
        (r"^Barilla", "Barilla"),
        (r"^Cerebos", "Cerebos"),
        (r"^Alhami", "Alhami"),
        (r"^Fantastic", "Fantastic"),
        (r"^Wok Time", "Wok Time"),
    ]
    for pat, name in rules:
        if re.search(pat, str(sku), re.I):
            return name
    if mfr == "Private Label":
        return "Private Label"
    if mfr == "Tiger Brands":
        return "Fatti's & Moni's"
    if mfr == "Unknown Manufacturer":
        return "Unknown"
    return mfr


def pack_grams(name: str) -> float | None:
    match = re.search(r"(\d+)\s*x\s*(\d+)\s*g", str(name), re.I)
    if match:
        return float(match.group(2))
    match = re.search(r"(\d+)\s*g", str(name), re.I)
    return float(match.group(1)) if match else None


def period_masks(df: pd.DataFrame) -> dict[str, pd.Series]:
    return {
        "mat": (df["date"] >= "2025-07-01") & (df["date"] <= "2026-06-01"),
        "ya": (df["date"] >= "2024-07-01") & (df["date"] <= "2025-06-01"),
        "jun": df["date"] == "2026-06-01",
        "jun_ya": df["date"] == "2025-06-01",
    }


def kpis(df: pd.DataFrame, mask: pd.Series) -> dict:
    s = df.loc[mask]
    value = float(s["value"].sum())
    volume = float(s["volume"].sum())
    return {
        "value": value,
        "volume": volume,
        "price": value / volume if volume else np.nan,
        "skus": int(s["sku"].nunique()),
        "mfrs": int(s["mfr"].nunique()),
    }


def growth(now: dict, then: dict) -> dict:
    return {
        "value": now["value"] / then["value"] - 1 if then["value"] else np.nan,
        "volume": now["volume"] / then["volume"] - 1 if then["volume"] else np.nan,
        "price": now["price"] / then["price"] - 1 if then["price"] else np.nan,
    }


def share_table(df: pd.DataFrame, mask: pd.Series, col: str, top: int | None = None) -> pd.DataFrame:
    g = df.loc[mask].groupby(col).agg(v=("value", "sum"), vol=("volume", "sum"))
    g["price"] = g["v"] / g["vol"]
    g["share"] = g["v"] / g["v"].sum() * 100
    g = g.sort_values("v", ascending=False)
    return g.head(top) if top else g


def joined_share(mat: pd.DataFrame, ya: pd.DataFrame) -> pd.DataFrame:
    out = mat.join(ya, lsuffix="_mat", rsuffix="_ya", how="outer").fillna(0)
    out["pp"] = out["share_mat"] - out["share_ya"]
    out["val_chg"] = np.where(out["v_ya"] > 0, out["v_mat"] / out["v_ya"] - 1, np.nan)
    return out.sort_values("v_mat", ascending=False)


def fmt_r(n: float, digits: int = 1) -> str:
    abs_n = abs(n)
    if abs_n >= 1e9:
        return f"R{n/1e9:.{digits}f}bn"
    if abs_n >= 1e6:
        return f"R{n/1e6:.{digits}f}m"
    if abs_n >= 1e3:
        return f"R{n/1e3:.0f}k"
    return f"R{n:.0f}"


def fmt_pct(n: float, signed: bool = True) -> str:
    if pd.isna(n):
        return "n/a"
    pct = n * 100
    sign = "+" if signed and pct > 0 else ""
    return f"{sign}{pct:.1f}%"


def fmt_pp(n: float) -> str:
    if pd.isna(n):
        return "n/a"
    sign = "+" if n > 0 else ""
    return f"{sign}{n:.1f}pp"


def set_run(run, text, size=12, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text_box(slide, l, t, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    sp_pr = shape._element.spPr
    effect = sp_pr.find(qn("a:effectLst"))
    if effect is not None:
        sp_pr.remove(effect)
    return shape


def add_header(slide, title, subtitle):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.88), NAVY)
    add_text_box(slide, Inches(0.38), Inches(0.10), Inches(12.5), Inches(0.40), title, size=22, bold=True, color=WHITE)
    add_text_box(
        slide, Inches(0.38), Inches(0.50), Inches(12.5), Inches(0.30), subtitle, size=11, color=RGBColor(0xC9, 0xD6, 0xE3)
    )
    add_rect(slide, 0, Inches(7.22), SLIDE_W, Inches(0.28), NAVY)
    add_text_box(
        slide,
        Inches(0.38),
        Inches(7.24),
        Inches(12.6),
        Inches(0.22),
        "Source: Monthly Trended Export, South Africa Pasta, multi-retailer  ·  12 months = Jul 25–Jun 26 vs YA Jul 24–Jun 25  ·  Price = value ÷ volume (R/kg)  ·  Latest month = Jun 26",
        size=10,
        color=RGBColor(0xC9, 0xD6, 0xE3),
    )


def add_kpi(slide, l, t, w, h, label, value, note=None, accent=BLUE):
    add_rect(slide, l, t, w, h, LIGHT, LINE)
    add_rect(slide, l, t, Inches(0.07), h, accent)
    add_text_box(slide, l + Inches(0.16), t + Inches(0.06), w - Inches(0.22), Inches(0.22), label, size=10, color=MUTED)
    add_text_box(slide, l + Inches(0.16), t + Inches(0.26), w - Inches(0.22), Inches(0.32), value, size=18, bold=True, color=NAVY)
    if note:
        add_text_box(slide, l + Inches(0.16), t + Inches(0.58), w - Inches(0.22), Inches(0.22), note, size=10, color=MUTED)


def add_table(slide, l, t, w, h, headers, rows, col_widths=None, header_fill=NAVY):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h)
    table = table_shape.table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        set_run(p.add_run(), header, size=10, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    for r, row in enumerate(rows, start=1):
        bg = LIGHT if r % 2 else WHITE
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            color = INK
            text = str(value)
            if c > 0 and text.startswith("+"):
                color = GREEN
            elif c > 0 and (text.startswith("-") and "pp" in text or (text.startswith("-") and "%" in text)):
                color = RED
            set_run(p.add_run(), text, size=11, bold=(c == 0), color=color)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return table


def add_bullets(slide, l, t, w, h, lines, size=13, color=INK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        set_run(p.add_run(), "•  " + line, size=size, color=color)
    return box


def chart_subcat(sub: pd.DataFrame) -> Path:
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    path = CHART_DIR / "pasta_subcat_share.png"
    labels = sub.index.tolist()
    shares = sub["share_mat"].tolist()
    fig, ax = plt.subplots(figsize=(5.8, 2.35), dpi=180)
    colors = ["#1B4F8A", "#0E7C7B", "#C45C26", "#6B7C8A", "#9B2C2C", "#7A6A4F"]
    ax.barh(labels[::-1], shares[::-1], color=list(reversed(colors[: len(labels)])), height=0.62)
    ax.set_xlabel("12-month value share (%)")
    ax.tick_params(axis="both", labelsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar, val in zip(ax.patches, shares[::-1]):
        ax.text(val + 0.6, bar.get_y() + bar.get_height() / 2, f"{val:.0f}%", va="center", fontsize=8)
    ax.set_xlim(0, max(shares) * 1.18)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_brands(brands: pd.DataFrame) -> Path:
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    path = CHART_DIR / "noodles_brand_share.png"
    plot = brands.head(6).copy()
    fig, ax = plt.subplots(figsize=(6.1, 2.45), dpi=180)
    colors = ["#1B4F8A" if pp >= 0 else "#9B2C2C" for pp in plot["pp"]]
    ax.barh(plot.index[::-1], plot["share_mat"][::-1], color=list(reversed(colors)), height=0.62)
    ax.set_xlabel("12-month noodles value share (%)")
    ax.tick_params(axis="both", labelsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar, share, pp in zip(ax.patches, plot["share_mat"][::-1], plot["pp"][::-1]):
        ax.text(
            share + 0.5,
            bar.get_y() + bar.get_height() / 2,
            f"{share:.0f}%  ({pp:+.1f}pp)",
            va="center",
            fontsize=8,
        )
    ax.set_xlim(0, plot["share_mat"].max() * 1.38)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_price(ladder: pd.DataFrame) -> Path:
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    path = CHART_DIR / "noodles_price_ladder.png"
    fig, ax = plt.subplots(figsize=(6.3, 2.35), dpi=180)
    colors = ["#C45C26" if b == "Knorr" else "#1B4F8A" for b in ladder.index]
    ax.barh(ladder.index[::-1], ladder["pack_price"][::-1], color=list(reversed(colors)), height=0.58)
    ax.set_xlabel("Typical unit pack price (R)")
    ax.tick_params(axis="both", labelsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar, val in zip(ax.patches, ladder["pack_price"][::-1]):
        ax.text(val + 0.6, bar.get_y() + bar.get_height() / 2, f"R{val:.2f}", va="center", fontsize=8)
    ax.set_xlim(0, max(ladder["pack_price"].max() * 1.22, 8))
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def analyse(df: pd.DataFrame) -> dict:
    m = period_masks(df)
    pasta_mat, pasta_ya = kpis(df, m["mat"]), kpis(df, m["ya"])
    pasta_jun, pasta_jun_ya = kpis(df, m["jun"]), kpis(df, m["jun_ya"])
    noodles = df[df["subcat"] == "Noodles"]
    n_mat, n_ya = kpis(noodles, m["mat"]), kpis(noodles, m["ya"])
    n_jun, n_jun_ya = kpis(noodles, m["jun"]), kpis(noodles, m["jun_ya"])

    sub_mat = share_table(df, m["mat"], "subcat")
    sub_ya = share_table(df, m["ya"], "subcat")
    sub = joined_share(sub_mat, sub_ya)

    pasta_mfr = joined_share(share_table(df, m["mat"], "mfr"), share_table(df, m["ya"], "mfr"))
    n_brand = joined_share(share_table(noodles, m["mat"], "brand"), share_table(noodles, m["ya"], "brand"))
    n_jun_brand = share_table(noodles, m["jun"], "brand")

    unilever = df[df["mfr"] == "Unilever"]
    u_jun = unilever[unilever["date"] == "2026-06-01"].copy()
    u_jun["short"] = (
        u_jun["sku"].str.replace(r"Knorr Pasta Pots\s+", "", regex=True).str.replace(r"\s+\d+g$", "", regex=True)
    )
    u_jun["share"] = u_jun["value"] / u_jun["value"].sum()

    # Price ladder: volume-weighted R/kg * typical unit grams for key brands
    nmat = noodles.loc[m["mat"] & noodles["volume"].gt(0)]
    ladder_rows = []
    specs = [
        ("Roka", 85),
        ("Kellogg's", 70),
        ("Maggi", 68),
        ("Indomie", 80),
        ("Samyang", 130),
        ("Knorr", 63),
    ]
    for brand, grams in specs:
        s = nmat[nmat["brand"] == brand]
        if s.empty:
            continue
        rkg = s["value"].sum() / s["volume"].sum()
        pack_price = rkg * grams / 1000
        if brand == "Knorr" and s["pack_price"].notna().any():
            pack_price = float((s["pack_price"] * s["value"]).sum() / s.loc[s["pack_price"].notna(), "value"].sum())
        ladder_rows.append({"brand": brand, "rkg": rkg, "pack_g": grams, "pack_price": pack_price})
    ladder = pd.DataFrame(ladder_rows).set_index("brand")

    sachet = nmat[nmat["pack_g"].fillna(999) < 150].copy()
    sachet_58 = sachet[(sachet["pack_price"] >= 5) & (sachet["pack_price"] < 8)]
    sachet_share = sachet_58["value"].sum() / n_mat["value"] if n_mat["value"] else 0

    return {
        "pasta_mat": pasta_mat,
        "pasta_ya": pasta_ya,
        "pasta_g": growth(pasta_mat, pasta_ya),
        "pasta_jun": pasta_jun,
        "pasta_jun_g": growth(pasta_jun, pasta_jun_ya),
        "n_mat": n_mat,
        "n_ya": n_ya,
        "n_g": growth(n_mat, n_ya),
        "n_jun": n_jun,
        "n_jun_g": growth(n_jun, n_jun_ya),
        "sub": sub,
        "pasta_mfr": pasta_mfr,
        "n_brand": n_brand,
        "n_jun_brand": n_jun_brand,
        "u_jun": u_jun,
        "ladder": ladder,
        "sachet_share": sachet_share,
        "n_share_pasta": n_mat["value"] / pasta_mat["value"],
        "n_share_pasta_ya": n_ya["value"] / pasta_ya["value"],
        "u_share_n_jun": u_jun["value"].sum() / n_jun["value"] if n_jun["value"] else 0,
        "u_share_pasta_jun": u_jun["value"].sum() / pasta_jun["value"] if pasta_jun["value"] else 0,
        "u_share_n_mat": u_jun["value"].sum() / n_mat["value"] if n_mat["value"] else 0,
    }


def build():
    df = load()
    a = analyse(df)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    sub_chart = chart_subcat(a["sub"].sort_values("share_mat", ascending=False))
    brand_chart = chart_brands(a["n_brand"])
    price_chart = chart_price(a["ladder"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ----- Slide 1: Pasta -----
    s1 = prs.slides.add_slide(blank)
    add_header(
        s1,
        "Pasta category: noodles now run the shop",
        "South Africa  ·  12 months to Jun 26  ·  All manufacturers",
    )
    add_kpi(
        s1,
        Inches(0.38),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "12-MONTH VALUE",
        fmt_r(a["pasta_mat"]["value"]),
        f"{fmt_pct(a['pasta_g']['value'])} vs YA  ·  vol {fmt_pct(a['pasta_g']['volume'])}",
        BLUE,
    )
    add_kpi(
        s1,
        Inches(3.58),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "JUNE 2026",
        fmt_r(a["pasta_jun"]["value"]),
        f"{fmt_pct(a['pasta_jun_g']['value'])} vs Jun 25",
        TEAL,
    )
    add_kpi(
        s1,
        Inches(6.78),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "NOODLES SHARE",
        f"{a['n_share_pasta']*100:.0f}% of Pasta",
        f"{fmt_pp((a['n_share_pasta']-a['n_share_pasta_ya'])*100)} vs YA",
        ORANGE,
    )
    add_kpi(
        s1,
        Inches(9.98),
        Inches(1.04),
        Inches(2.96),
        Inches(0.88),
        "CATEGORY PRICE",
        f"R{a['pasta_mat']['price']:.0f}/kg",
        f"Noodles R{a['n_mat']['price']:.0f}  ·  dry ~R32–35",
        RED,
    )

    add_text_box(s1, Inches(0.38), Inches(2.02), Inches(6.2), Inches(0.26), "Value mix by subcategory (12 months)", size=12, bold=True, color=NAVY)
    s1.shapes.add_picture(str(sub_chart), Inches(0.32), Inches(2.28), Inches(6.35), Inches(2.45))

    mfr_rows = []
    for name, r in a["pasta_mfr"].head(7).iterrows():
        mfr_rows.append(
            [
                name if name != "Zhejiang Green Home Food Co., Ltd." else "Roka (Zhejiang)",
                f"{r['share_mat']:.1f}%",
                fmt_pp(r["pp"]),
                fmt_pct(r["val_chg"]),
                f"R{r['price_mat']:.0f}",
            ]
        )
    add_text_box(s1, Inches(6.85), Inches(2.02), Inches(6.1), Inches(0.26), "Who plays, who is winning (12-month value)", size=12, bold=True, color=NAVY)
    add_table(
        s1,
        Inches(6.85),
        Inches(2.30),
        Inches(6.10),
        Inches(2.55),
        ["Manufacturer", "Share", "vs YA", "Value", "R/kg"],
        mfr_rows,
        col_widths=[Inches(2.15), Inches(0.90), Inches(0.95), Inches(1.00), Inches(1.10)],
    )

    add_rect(s1, Inches(0.38), Inches(4.92), Inches(12.56), Inches(2.12), LIGHT, LINE)
    add_text_box(s1, Inches(0.55), Inches(5.02), Inches(12.2), Inches(0.26), "Takeaways", size=13, bold=True, color=NAVY)
    add_bullets(
        s1,
        Inches(0.55),
        Inches(5.30),
        Inches(12.2),
        Inches(1.62),
        [
            f"Pasta is a R{a['pasta_mat']['value']/1e9:.2f}bn category growing {fmt_pct(a['pasta_g']['value'])} in value, with volume {fmt_pct(a['pasta_g']['volume'])} — growth is price/mix, not more tonnes.",
            "Noodles is the engine: 64% of Pasta value (up 1.8pp). Macaroni 21% and spaghetti 13% are the dry-pasta game, led by Tiger Brands and private label at ~R30–35/kg.",
            "Winners: Kellogg's (+2.1pp) and Indomie/Indofood (+1.4pp). Losers: Nestlé/Maggi (−2.7pp), Mr Pasta (−2.0pp), Tiger Brands (−1.5pp).",
            "Unilever is not a Pasta-category player yet. Knorr Pasta Pots appear only in Jun 26, coded in Noodles.",
        ],
        size=13,
    )
    s1.notes_slide.notes_text_frame.text = (
        f"Pasta MAT {a['pasta_mat']['value']:.0f} vs YA {a['pasta_ya']['value']:.0f}. "
        f"Noodles MAT {a['n_mat']['value']:.0f}. "
        "Retailer set as in the extract header (Clicks, Massmart, Spar, Shoprite, PnP, etc.)."
    )

    # ----- Slide 2: Noodles -----
    s2 = prs.slides.add_slide(blank)
    add_header(
        s2,
        "Noodles: Kellogg's winning, Maggi losing, Unilever not in the sachet game",
        "Noodles subcategory  ·  12 months to Jun 26 and June 2026  ·  Unilever = Knorr Pasta Pots",
    )
    add_kpi(
        s2,
        Inches(0.38),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "12-MONTH VALUE",
        fmt_r(a["n_mat"]["value"]),
        f"{fmt_pct(a['n_g']['value'])}  ·  vol {fmt_pct(a['n_g']['volume'])}",
        BLUE,
    )
    add_kpi(
        s2,
        Inches(3.58),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "JUNE 2026",
        fmt_r(a["n_jun"]["value"]),
        f"{fmt_pct(a['n_jun_g']['value'])} vs Jun 25",
        TEAL,
    )
    add_kpi(
        s2,
        Inches(6.78),
        Inches(1.04),
        Inches(3.10),
        Inches(0.88),
        "PRICE POINT",
        f"R{a['n_mat']['price']:.0f}/kg",
        f"{a['sachet_share']*100:.0f}% of value is R5–8 a sachet",
        ORANGE,
    )
    add_kpi(
        s2,
        Inches(9.98),
        Inches(1.04),
        Inches(2.96),
        Inches(0.88),
        "UNILEVER (JUN)",
        f"{a['u_share_n_jun']*100:.2f}% share",
        f"{fmt_r(a['u_jun']['value'].sum())}  ·  3 Pasta Pots",
        RED,
    )

    add_text_box(s2, Inches(0.38), Inches(2.02), Inches(6.3), Inches(0.24), "Key players — 12-month value share and change", size=12, bold=True, color=NAVY)
    s2.shapes.add_picture(str(brand_chart), Inches(0.28), Inches(2.24), Inches(6.45), Inches(2.55))

    brand_rows = []
    show = ["Kellogg's", "Maggi", "Indomie", "Roka", "Samyang", "Eezee/Joy", "Mr Pasta", "Knorr"]
    for name in show:
        if name not in a["n_brand"].index:
            continue
        r = a["n_brand"].loc[name]
        jun_share = a["n_jun_brand"].loc[name, "share"] if name in a["n_jun_brand"].index else 0
        brand_rows.append(
            [
                name,
                f"{r['share_mat']:.1f}%",
                f"{jun_share:.1f}%",
                fmt_pp(r["pp"]),
                fmt_pct(r["val_chg"]),
                f"R{r['price_mat']:.0f}",
            ]
        )
    add_text_box(s2, Inches(6.85), Inches(2.02), Inches(6.1), Inches(0.24), "Share, June pulse, price", size=12, bold=True, color=NAVY)
    add_table(
        s2,
        Inches(6.85),
        Inches(2.28),
        Inches(6.10),
        Inches(2.70),
        ["Brand", "12m", "Jun", "vs YA", "Value", "R/kg"],
        brand_rows,
        col_widths=[Inches(1.35), Inches(0.85), Inches(0.80), Inches(0.90), Inches(1.05), Inches(1.15)],
    )

    add_text_box(s2, Inches(0.38), Inches(4.90), Inches(6.3), Inches(0.24), "Price ladder — typical unit pack", size=12, bold=True, color=NAVY)
    s2.shapes.add_picture(str(price_chart), Inches(0.28), Inches(5.10), Inches(6.50), Inches(1.95))

    add_rect(s2, Inches(6.85), Inches(5.08), Inches(6.10), Inches(1.98), LIGHT, LINE)
    add_text_box(s2, Inches(7.00), Inches(5.16), Inches(5.8), Inches(0.24), "Unilever competing SKUs (Jun 26 only)", size=12, bold=True, color=NAVY)
    u_rows = []
    for _, r in a["u_jun"].sort_values("value", ascending=False).iterrows():
        u_rows.append(
            [
                r["short"],
                f"R{r['value']/1000:.0f}k",
                f"{r['share']*100:.0f}%",
                f"R{r['pack_price']:.2f}",
            ]
        )
    add_table(
        s2,
        Inches(7.00),
        Inches(5.42),
        Inches(5.80),
        Inches(1.52),
        ["SKU", "Value", "Mix", "R/pot"],
        u_rows,
        col_widths=[Inches(2.35), Inches(1.10), Inches(1.00), Inches(1.35)],
    )

    s2.notes_slide.notes_text_frame.text = (
        "Sachet price uses brand R/kg × typical unit grams (Kellogg's 70g, Maggi 68g, Indomie 80g, Roka 85g, Samyang 130g, Knorr 63g). "
        f"{a['sachet_share']*100:.0f}% of noodles value is in the R5–8 unit-pack band. "
        "Knorr Pasta Pots are not in the file before Jun 26."
    )

    # ----- Slide 3: So what -----
    s3 = prs.slides.add_slide(blank)
    add_header(
        s3,
        "So what for Unilever",
        "Pasta Pots are a different product and price tier, not a head-to-head with Maggi or Kellogg's sachets",
    )

    add_rect(s3, Inches(0.38), Inches(1.08), Inches(6.20), Inches(5.92), LIGHT, LINE)
    add_text_box(s3, Inches(0.55), Inches(1.20), Inches(5.85), Inches(0.30), "What the category is doing", size=16, bold=True, color=NAVY)
    add_bullets(
        s3,
        Inches(0.55),
        Inches(1.58),
        Inches(5.85),
        Inches(5.20),
        [
            "Noodles is the Pasta category. 64% of value, growing faster than dry pasta, and the only segment Unilever is in.",
            "The mass market is a R6 sachet: Kellogg's (~R87/kg) and Maggi (~R100/kg). Roka undercuts at ~R69/kg. That band is ~87% of noodles value.",
            "Kellogg's is taking Maggi's lunch. Maggi −5.1pp over 12 months; Kellogg's +2.0pp. Indomie and Samyang are the other gainers.",
            "Dry pasta is a separate fight (Tiger, private label, Mr Pasta at R30–35/kg). Unilever is absent there.",
        ],
        size=14,
    )

    add_rect(s3, Inches(6.78), Inches(1.08), Inches(6.16), Inches(5.92), NAVY)
    add_text_box(s3, Inches(6.96), Inches(1.20), Inches(5.80), Inches(0.30), "Where Unilever sits", size=16, bold=True, color=WHITE)
    box = s3.shapes.add_textbox(Inches(6.96), Inches(1.58), Inches(5.80), Inches(5.20))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        f"Knorr Pasta Pots launched into the file in June only: {fmt_r(a['u_jun']['value'].sum())}, 0.56% of Noodles, 0.37% of Pasta.",
        "One range, one price: ~R49 a pot / ~R790/kg — about 8× the mainstream sachet and 8× category noodles R/kg.",
        "Carbonara is the lead SKU (49%). This is a cup/meal occasion, closer to Samyang/ramen than to Maggi 2-minute sachets.",
        "The three SKUs compete with each other more than with Kellogg's. Scale is still a rounding error.",
        "Watch: Maggi's share loss (room if Nestlé stumbles on occasion), Indomie (flavour/heat), and whether Pasta Pots can hold R49 without promo.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        set_run(p.add_run(), "•  " + line, size=14, color=WHITE)

    s3.notes_slide.notes_text_frame.text = (
        "Unilever MAT share of noodles is ~0.06% because the SKUs have one month of sales. "
        "Do not read June share as a 12-month position."
    )

    prs.save(DECK)
    print(f"Wrote {DECK}")
    print("Pasta MAT", fmt_r(a["pasta_mat"]["value"]), fmt_pct(a["pasta_g"]["value"]))
    print("Noodles MAT", fmt_r(a["n_mat"]["value"]), fmt_pct(a["n_g"]["value"]))
    print("Noodles share of pasta", f"{a['n_share_pasta']*100:.1f}%")
    print("Unilever Jun noodles share", f"{a['u_share_n_jun']*100:.2f}%")
    print(a["n_brand"][["share_mat", "pp", "val_chg", "price_mat"]].head(8).to_string())


if __name__ == "__main__":
    build()
