#!/usr/bin/env python3
"""Build the Pasta / Noodles review in the Smollan Unilever template style."""

from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from build_deck import analyse, fmt_pct, fmt_pp, fmt_r, load

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PDF = ROOT / "templates" / "Smollan_Unilever_Template.pdf"
ASSETS = ROOT / "output" / "template_assets"
CHART_DIR = ROOT / "output" / "charts"
DECK = ROOT / "output" / "Unilever_Pasta_Noodles_Jun2026.pptx"

NAVY = RGBColor(0x00, 0x30, 0x58)
BLUE = RGBColor(0x00, 0x78, 0xC0)
CYAN = RGBColor(0x00, 0xB8, 0xF0)
LIME = RGBColor(0x88, 0xC0, 0x40)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1F, 0x7A, 0x4D)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x6B, 0x73)
LIGHT = RGBColor(0xEE, 0xF3, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xE6, 0xF1, 0xF8)
FOOTER = RGBColor(0x6A, 0x7A, 0x86)
ACTION_GREEN = RGBColor(0xD9, 0xEF, 0xC7)
ACTION_BLUE = RGBColor(0xD6, 0xEB, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def extract_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    doc = pymupdf.open(TEMPLATE_PDF)
    mapping = {
        "title_bg": (0, 11),
        "logo_white_raw": (0, 13),
        "agenda_bg": (1, 24),
        "logo_content_raw": (5, 39),
        "thanks_bg": (12, 125),
    }
    paths = {}
    for name, (page_i, xref) in mapping.items():
        pix = pymupdf.Pixmap(doc, xref)
        if pix.n - pix.alpha > 3:
            pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
        path = ASSETS / f"{name}.png"
        pix.save(str(path))
        paths[name] = path

    def knock_and_crop(src: Path, dest: Path, thresh: int = 28) -> Path:
        im = Image.open(src).convert("RGBA")
        arr = np.array(im)
        r, g, b = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2]
        mask = (r.astype(int) + g.astype(int) + b.astype(int)) < thresh * 3
        arr[:, :, 3] = np.where(mask, 0, 255)
        im = Image.fromarray(arr)
        ys, xs = np.where(arr[:, :, 3] > 10)
        crop = im.crop((int(xs.min()), int(ys.min()), int(xs.max()) + 1, int(ys.max()) + 1))
        crop.save(dest)
        return dest

    paths["logo_white"] = knock_and_crop(paths["logo_white_raw"], ASSETS / "logo_white.png")
    paths["logo_content"] = knock_and_crop(paths["logo_content_raw"], ASSETS / "logo_content.png")
    return paths


def set_run(run, text, size=12, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(slide, l, t, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    sp_pr = shape._element.spPr
    effect = sp_pr.find(qn("a:effectLst"))
    if effect is not None:
        sp_pr.remove(effect)
    if rounded:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_picture_bg(slide, path: Path):
    slide.shapes.add_picture(str(path), 0, 0, SLIDE_W, SLIDE_H)


def content_chrome(slide, assets, title, insight, page_no):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT)
    slide.shapes.add_picture(str(assets["logo_content"]), Inches(0.28), Inches(0.16), Inches(1.85), Inches(0.52))
    add_text(slide, Inches(2.30), Inches(0.18), Inches(10.7), Inches(0.42), title, size=22, bold=True, color=INK)
    add_rect(slide, Inches(0.28), Inches(0.78), Inches(12.76), Inches(0.72), WHITE, RGBColor(0xD5, 0xDE, 0xE6), rounded=True)
    add_text(slide, Inches(0.42), Inches(0.86), Inches(12.48), Inches(0.56), insight, size=13, color=INK)
    add_text(slide, Inches(0.28), Inches(7.22), Inches(4.5), Inches(0.22), "A SMOLLAN COMPANY", size=9, color=FOOTER)
    add_text(
        slide,
        Inches(5.4),
        Inches(7.22),
        Inches(7.65),
        Inches(0.22),
        f"EMPOWERING RETAIL PRECISION WITH DATA, SOFTWARE & AI    {page_no}",
        size=9,
        color=FOOTER,
        align=PP_ALIGN.RIGHT,
    )


def add_table(slide, l, t, w, h, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h)
    table = table_shape.table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        set_run(p.add_run(), header, size=10, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        bg = ROW_ALT if r % 2 else WHITE
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            text = str(value)
            color = INK
            if c > 0 and text.startswith("+"):
                color = GREEN
            elif c > 0 and text.startswith("-"):
                color = RED
            set_run(p.add_run(), text, size=11, bold=(c == 0), color=color)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return table


def style_charts():
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "axes.edgecolor": "#D5DEE6",
            "axes.labelcolor": "#1A1A1A",
            "xtick.color": "#1A1A1A",
            "ytick.color": "#1A1A1A",
            "text.color": "#1A1A1A",
        }
    )


def doughnut(share: float, path: Path, label: str):
    fig, ax = plt.subplots(figsize=(2.4, 2.4), dpi=160)
    ax.pie(
        [share, 1 - share],
        colors=["#003058", "#D5DEE6"],
        startangle=90,
        wedgeprops={"width": 0.42, "edgecolor": "white"},
    )
    ax.text(0, 0.08, f"{share*100:.0f}%", ha="center", va="center", fontsize=16, fontweight="bold", color="#003058")
    ax.text(0, -0.18, label, ha="center", va="center", fontsize=7, color="#5C6B73")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white", transparent=True)
    plt.close(fig)


def barh_chart(labels, values, path: Path, xlabel: str, colors=None, value_fmt=lambda v: f"{v:.0f}%"):
    fig, ax = plt.subplots(figsize=(5.6, 2.35), dpi=170)
    cols = colors or ["#0078C0"] * len(labels)
    ax.barh(list(reversed(labels)), list(reversed(values)), color=list(reversed(cols)), height=0.58)
    ax.set_xlabel(xlabel, fontsize=8)
    ax.tick_params(labelsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    xmax = max(values) * 1.28 if values else 1
    ax.set_xlim(0, xmax)
    for bar, val in zip(ax.patches, reversed(values)):
        ax.text(val + xmax * 0.02, bar.get_y() + bar.get_height() / 2, value_fmt(val), va="center", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def kpi_panel(slide, l, t, w, h, header, big, body, chart_path=None):
    add_rect(slide, l, t, w, h, WHITE, RGBColor(0xD5, 0xDE, 0xE6), rounded=True)
    add_rect(slide, l, t, w, Inches(0.36), NAVY)
    add_text(slide, l + Inches(0.12), t + Inches(0.04), w - Inches(0.2), Inches(0.28), header, size=12, bold=True, color=WHITE)
    add_text(slide, l + Inches(0.14), t + Inches(0.44), w - Inches(0.28), Inches(0.38), big, size=24, bold=True, color=NAVY)
    body_w = w - Inches(2.05) if chart_path else w - Inches(0.28)
    add_text(slide, l + Inches(0.14), t + Inches(0.86), body_w, Inches(1.50), body, size=12, color=INK)
    if chart_path:
        slide.shapes.add_picture(str(chart_path), l + w - Inches(1.90), t + Inches(0.48), Inches(1.72), Inches(1.72))


def build():
    assets = extract_assets()
    df = load()
    a = analyse(df)
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    style_charts()

    doughnut(a["n_share_pasta"], CHART_DIR / "sm_noodles_share.png", "of Pasta value")
    doughnut(a["sachet_share"], CHART_DIR / "sm_sachet_share.png", "R5–8 sachets")

    sub = a["sub"].sort_values("share_mat", ascending=False)
    barh_chart(
        sub.index.tolist(),
        sub["share_mat"].tolist(),
        CHART_DIR / "sm_subcat.png",
        "12-month value share (%)",
        colors=["#003058", "#0078C0", "#00B8F0", "#88C040", "#5C6B73", "#C0392B"],
    )
    brands = a["n_brand"].head(6)
    barh_chart(
        brands.index.tolist(),
        brands["share_mat"].tolist(),
        CHART_DIR / "sm_brands.png",
        "12-month noodles value share (%)",
        colors=["#0078C0" if pp >= 0 else "#C0392B" for pp in brands["pp"]],
        value_fmt=lambda v: f"{v:.0f}%",
    )
    ladder = a["ladder"]
    barh_chart(
        ladder.index.tolist(),
        ladder["pack_price"].tolist(),
        CHART_DIR / "sm_price.png",
        "Typical unit pack price (R)",
        colors=["#00B8F0" if b != "Knorr" else "#88C040" for b in ladder.index],
        value_fmt=lambda v: f"R{v:.2f}",
    )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    add_picture_bg(s, assets["title_bg"])
    s.shapes.add_picture(str(assets["logo_white"]), Inches(0.42), Inches(0.42), Inches(2.55), Inches(0.78))
    add_text(s, Inches(0.45), Inches(2.35), Inches(6.4), Inches(0.85), "Unilever Presentation", size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.45), Inches(3.20), Inches(6.4), Inches(0.40), "Pasta & Noodles Category Review", size=18, color=WHITE)
    add_text(s, Inches(0.45), Inches(3.62), Inches(6.4), Inches(0.32), "South Africa  |  12 months to June 2026", size=14, color=CYAN)
    add_text(s, Inches(0.45), Inches(6.55), Inches(5.5), Inches(0.28), "Updated: 18/08/2026", size=13, bold=True, color=LIME)
    add_text(s, Inches(0.45), Inches(7.12), Inches(5.5), Inches(0.22), "A SMOLLAN COMPANY", size=10, color=WHITE)

    # 2. Agenda
    s = prs.slides.add_slide(blank)
    add_picture_bg(s, assets["agenda_bg"])
    s.shapes.add_picture(str(assets["logo_white"]), Inches(0.42), Inches(0.28), Inches(2.15), Inches(0.64))
    add_text(s, Inches(0.45), Inches(1.10), Inches(7.6), Inches(0.42), "AGENDA -", size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.45), Inches(1.52), Inches(7.6), Inches(0.32), "Key business questions from the brief", size=16, color=WHITE)
    add_text(s, Inches(0.45), Inches(1.90), Inches(7.6), Inches(0.28), "Primary analysis: 12 months Jul 25–Jun 26  ·  latest month Jun 26", size=13, color=RGBColor(0xD6, 0xEE, 0xF8))
    items = [
        ("01  —  Category growth", "Where is the Pasta value / volume opportunity?"),
        ("02  —  Key players", "Who plays in Pasta and Noodles, and who is winning or losing?"),
        ("03  —  Pack & price", "What is the Noodles price architecture?"),
        ("04  —  Unilever position", "Where do Knorr Pasta Pots sit versus competing SKUs?"),
        ("05  —  Commercial action", "What should Unilever Protect, Watch and Test?"),
    ]
    box = s.shapes.add_textbox(Inches(0.45), Inches(2.35), Inches(7.7), Inches(4.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (title, q) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        set_run(p.add_run(), title, size=16, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(12)
        set_run(p2.add_run(), q, size=13, color=RGBColor(0xD6, 0xEE, 0xF8))
    add_text(s, Inches(0.45), Inches(7.12), Inches(5.5), Inches(0.22), "A SMOLLAN COMPANY", size=10, color=WHITE)

    # 3. Pasta insights
    s = prs.slides.add_slide(blank)
    content_chrome(
        s,
        assets,
        "Pasta category - Key market insights",
        "Pasta is a R3.3bn category growing +5.2% in value with volume slightly down. Noodles now contribute 64% of value and are the engine of growth.",
        3,
    )
    kpi_panel(
        s,
        Inches(0.28),
        Inches(1.64),
        Inches(6.28),
        Inches(2.55),
        "Category growth",
        fmt_r(a["pasta_mat"]["value"]),
        f"12-month value {fmt_pct(a['pasta_g']['value'])} vs YA. Volume {fmt_pct(a['pasta_g']['volume'])}. June 2026 {fmt_r(a['pasta_jun']['value'])}, {fmt_pct(a['pasta_jun_g']['value'])} vs Jun 25. Growth is price/mix, not more tonnes.",
    )
    kpi_panel(
        s,
        Inches(6.72),
        Inches(1.64),
        Inches(6.32),
        Inches(2.55),
        "Noodles concentration",
        f"{a['n_share_pasta']*100:.0f}%",
        f"Noodles is 64.4% of Pasta value, up 1.8pp vs year ago. Macaroni 21% and spaghetti 13% are dry pasta. Unilever is only in Noodles.",
        CHART_DIR / "sm_noodles_share.png",
    )
    kpi_panel(
        s,
        Inches(0.28),
        Inches(4.32),
        Inches(6.28),
        Inches(2.70),
        "Who is winning / losing",
        "Kellogg's +2.1pp",
        "Kellogg's 29% of Pasta (+2.1pp) and Indomie +1.4pp. Nestlé −2.7pp, Mr Pasta −2.0pp, Tiger −1.5pp. Tiger leads dry pasta; Kellogg's/Nestlé lead because Noodles is most of the value.",
    )
    kpi_panel(
        s,
        Inches(6.72),
        Inches(4.32),
        Inches(6.32),
        Inches(2.70),
        "Price architecture",
        f"R{a['pasta_mat']['price']:.0f}/kg",
        f"Category average R{a['pasta_mat']['price']:.0f}/kg. Noodles R{a['n_mat']['price']:.0f}/kg. Dry pasta (macaroni/spaghetti) ~R32–35/kg. Two different businesses sit inside one Pasta file.",
    )

    # 4. Noodles players
    s = prs.slides.add_slide(blank)
    content_chrome(
        s,
        assets,
        "Noodles - Who plays, who is winning, what price point?",
        "Kellogg's is taking share from Maggi. The mass market is a R5–8 sachet. Knorr Pasta Pots entered in June at ~R49 — a different occasion, not a sachet competitor.",
        4,
    )
    add_text(s, Inches(0.32), Inches(1.62), Inches(6.2), Inches(0.26), "12-month noodles value share", size=12, bold=True, color=NAVY)
    s.shapes.add_picture(str(CHART_DIR / "sm_brands.png"), Inches(0.28), Inches(1.88), Inches(6.25), Inches(2.45))

    rows = []
    for name in ["Kellogg's", "Maggi", "Indomie", "Roka", "Samyang", "Eezee/Joy", "Mr Pasta", "Knorr"]:
        if name not in a["n_brand"].index:
            continue
        r = a["n_brand"].loc[name]
        jun = a["n_jun_brand"].loc[name, "share"] if name in a["n_jun_brand"].index else 0
        rows.append(
            [name, f"{r['share_mat']:.1f}%", f"{jun:.1f}%", fmt_pp(r["pp"]), fmt_pct(r["val_chg"]), f"R{r['price_mat']:.0f}"]
        )
    add_text(s, Inches(6.70), Inches(1.62), Inches(6.3), Inches(0.26), "Competitive scorecard", size=12, bold=True, color=NAVY)
    add_table(
        s,
        Inches(6.70),
        Inches(1.90),
        Inches(6.32),
        Inches(2.55),
        ["Brand", "12m", "Jun 26", "vs YA", "Value", "R/kg"],
        rows,
        col_widths=[Inches(1.35), Inches(0.85), Inches(0.95), Inches(0.95), Inches(1.05), Inches(1.17)],
    )

    add_text(s, Inches(0.32), Inches(4.42), Inches(6.2), Inches(0.24), "Price ladder — typical unit pack", size=12, bold=True, color=NAVY)
    s.shapes.add_picture(str(CHART_DIR / "sm_price.png"), Inches(0.22), Inches(4.64), Inches(6.35), Inches(2.35))

    add_rect(s, Inches(6.70), Inches(4.55), Inches(6.32), Inches(2.48), WHITE, RGBColor(0xD5, 0xDE, 0xE6), rounded=True)
    add_rect(s, Inches(6.70), Inches(4.55), Inches(6.32), Inches(0.34), NAVY)
    add_text(s, Inches(6.84), Inches(4.58), Inches(6.0), Inches(0.28), "Unilever competing SKUs — June 2026 only", size=12, bold=True, color=WHITE)
    u_rows = []
    for _, r in a["u_jun"].sort_values("value", ascending=False).iterrows():
        u_rows.append([r["short"], f"R{r['value']/1000:.0f}k", f"{r['share']*100:.0f}%", f"R{r['pack_price']:.2f}", f"R{r['price']:.0f}/kg"])
    add_table(
        s,
        Inches(6.82),
        Inches(4.96),
        Inches(6.08),
        Inches(1.95),
        ["SKU", "Value", "Mix", "R/pot", "R/kg"],
        u_rows,
        col_widths=[Inches(2.05), Inches(0.95), Inches(0.80), Inches(1.05), Inches(1.23)],
    )

    # 5. Action
    s = prs.slides.add_slide(blank)
    content_chrome(
        s,
        assets,
        "Noodles - How Unilever should read this category",
        f"Knorr Pasta Pots are a R{a['u_jun']['value'].sum()/1e6:.2f}m June entry (0.56% of Noodles). They do not compete in the R6 sachet game that is {a['sachet_share']*100:.0f}% of noodles value. Treat them as a cup/meal occasion and keep the three SKUs from fragmenting.",
        5,
    )
    cols = [
        (
            "Protect the read",
            BLUE,
            ACTION_BLUE,
            [
                "Noodles is the Pasta category (64% of value, +8.2%).",
                "Kellogg's 46% and Maggi 29% set the competitive frame.",
                "Dry pasta is Tiger / private label at R32–35/kg — Unilever is absent.",
            ],
            "Do not brief Pasta Pots as a Maggi 2-minute sachet competitor.",
        ),
        (
            "Watch Maggi's loss",
            BLUE,
            ACTION_BLUE,
            [
                "Maggi −5.1pp over 12 months; Kellogg's +2.0pp.",
                "Indomie +2.1pp and Samyang +1.8pp are the other gainers.",
                "Mr Pasta −2.3pp — value/mid is under pressure.",
            ],
            "Use Maggi's share loss as the main competitive opening to watch, not to copy on price.",
        ),
        (
            "Price architecture",
            RGBColor(0xC4, 0x5C, 0x26),
            ACTION_BLUE,
            [
                "Roka ~R5.88  ·  Kellogg's ~R6.09  ·  Maggi ~R6.79.",
                "Indomie ~R7.90. Samyang ~R39 (ramen/cup).",
                "Knorr Pasta Pots ~R49 / ~R790/kg — ~8× sachet.",
            ],
            "Hold ~R49. Promo into the R6 band would destroy the proposition.",
        ),
        (
            "Test the range",
            BLUE,
            ACTION_GREEN,
            [
                "Carbonara is 49% of the line; Mushroom is the smallest.",
                "Three SKUs at one price compete with each other first.",
                "June only — no 12-month Unilever trend yet.",
            ],
            "Keep distribution and visibility on Carbonara; measure incrementality vs cannibalisation across the three pots.",
        ),
    ]
    x = Inches(0.28)
    width = Inches(3.12)
    gap = Inches(0.14)
    for title, head_color, action_fill, bullets, action in cols:
        add_rect(s, x, Inches(1.62), width, Inches(5.40), WHITE, RGBColor(0xD5, 0xDE, 0xE6), rounded=True)
        add_text(s, x + Inches(0.12), Inches(1.70), width - Inches(0.22), Inches(0.55), title, size=14, bold=True, color=head_color)
        add_text(s, x + Inches(0.12), Inches(2.22), width - Inches(0.22), Inches(0.24), "Why?", size=11, bold=True, color=NAVY)
        box = s.shapes.add_textbox(x + Inches(0.12), Inches(2.46), width - Inches(0.22), Inches(2.45))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            set_run(p.add_run(), "•  " + line, size=11, color=INK)
        add_rect(s, x + Inches(0.10), Inches(5.28), width - Inches(0.20), Inches(1.58), action_fill, rounded=True)
        add_text(s, x + Inches(0.20), Inches(5.34), width - Inches(0.38), Inches(0.22), "Action", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.20), Inches(5.56), width - Inches(0.38), Inches(1.22), action, size=11, color=INK)
        x = Emu(x + width + gap)

    # 6. Thank you
    s = prs.slides.add_slide(blank)
    add_picture_bg(s, assets["thanks_bg"])
    s.shapes.add_picture(str(assets["logo_white"]), Inches(0.38), Inches(0.28), Inches(2.05), Inches(0.62))
    add_text(s, Inches(7.15), Inches(2.85), Inches(5.6), Inches(0.70), "Thank You!", size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s, Inches(7.15), Inches(3.58), Inches(2.8), Inches(0.03), RGBColor(0x7A, 0x9A, 0xB8))
    add_text(s, Inches(7.15), Inches(3.72), Inches(5.6), Inches(0.40), "Joseph Hlongwane", size=18, bold=True, color=LIME)
    add_text(s, Inches(0.40), Inches(7.12), Inches(5.5), Inches(0.22), "A SMOLLAN COMPANY", size=10, color=WHITE)

    prs.save(DECK)
    print(f"Wrote {DECK}")
    print("slides", len(prs.slides))


if __name__ == "__main__":
    build()
